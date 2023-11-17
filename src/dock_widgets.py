from __future__ import annotations

import csv
import os
import random
import webbrowser
from typing import TYPE_CHECKING

import pptx
from PIL import Image
from PySide6.QtCore import Qt, Signal
from PySide6.QtWidgets import QWidget, QDialog, QApplication, QListWidgetItem, QListWidget

from src import document_builder
from src.basic_config import base_path
from src.database import db
from src.datatypes import Regeltest, RegeltestIcon, SelfTestMode
from src.regeltestcreator import RegeltestSetup, RegeltestSaveDialog
from src.ui_regeltest_creator_dockwidget import Ui_regeltest_creator_dockwidget
from src.ui_self_test_dockwidget import Ui_self_test_dockwidget

if TYPE_CHECKING:
    from src.main_application import MainWindow


def generate_powerpoint_questionsets(questions_a: List[Question]):
    # Make a copy of the list to avoid modifying the original
    questions_b = list(questions_a)

    # Iterate over the shuffled list in pairs
    invalid = True
    matched_a, matched_b = [], []

    while invalid:
        # Shuffle the list randomly
        random.shuffle(questions_b)

        for i in range(len(questions_a)):
            # Check if the pair contains two identical questions_a
            if questions_a[i].question == questions_b[i].question:
                # If it does, swap one of them with the next element in the list until a valid pair is found
                for j in range(i + 1, len(questions_b)):
                    if questions_a[i].question != questions_b[j].question:
                        questions_b[i], questions_b[j] = questions_b[j], questions_b[i]
                        break

            # Add the valid pair to the list of pairs
            matched_a += [questions_a[i]]
            matched_b += [questions_b[i]]

        # small validity check
        assert len(matched_a) == len(questions_a)
        for i in range(len(matched_a)):
            if matched_a[0] != matched_b[1]:
                invalid = False
            else:
                invalid = True
                matched_a, matched_b = [], []
                break

    return matched_a, matched_b


class RegeltestCreatorDockwidget(QWidget, Ui_regeltest_creator_dockwidget):
    def __init__(self, main_window: MainWindow):
        super(RegeltestCreatorDockwidget, self).__init__(main_window)
        self.ui = Ui_regeltest_creator_dockwidget()
        self.ui.setupUi(self)

        self.ui.regeltest_list.setAcceptDrops(True)
        self.ui.regeltest_list.model().rowsInserted.connect(self.regeltest_list_updated)
        self.ui.regeltest_list.model().rowsRemoved.connect(self.regeltest_list_updated)

        self.ui.add_questionlist.clicked.connect(self.setup_regeltest)
        self.ui.clear_questionlist.clicked.connect(self.clear_questionlist)

        self.ui.create_regeltest.clicked.connect(self.create_regeltest)

    def clear_questionlist(self):
        self.ui.regeltest_list.clear()
        self.ui.regeltest_list.questions.clear()
        self.regeltest_list_updated()

    def regeltest_list_updated(self):
        self.ui.regeltest_stats.setText(
            f"{self.ui.regeltest_list.count()} Fragen selektiert ({self.ui.regeltest_list.count() * 2} Punkte)")

    def setup_regeltest(self):
        regeltest_setup = RegeltestSetup(self)
        if regeltest_setup.exec():
            for question in regeltest_setup.collect_questions():
                self.ui.regeltest_list.add_question(question)

    def create_regeltest(self):
        questions = []
        for signature in self.ui.regeltest_list.questions:
            questions += [db.get_question(signature)]
        settings = RegeltestSaveDialog(questions, self)
        settings.ui.title_edit.setFocus()
        result = settings.exec()
        pdf_path = settings.ui.pdf_edit.text()
        ppt_path = settings.ui.ppt_edit.text()
        csv_path = settings.ui.csv_edit.text()
        archive_regeltest = settings.ui.regeltest_archive_checkBox.isChecked()
        if result == QDialog.Accepted:
            selected_questions = settings.get_questions()
            QApplication.setOverrideCursor(Qt.WaitCursor)
            if settings.ui.icon_path_edit.text():
                icon = Image.open(settings.ui.icon_path_edit.text())
                icon_db = db.get_or_create(RegeltestIcon, icon=icon.tobytes())
            else:
                icon = None
                icon_db = None
            if (pdf_path or csv_path or ppt_path) and archive_regeltest:
                regeltest = Regeltest(title=settings.ui.title_edit.text(), icon=icon_db,
                                      selected_questions=selected_questions)
                db.add_object(regeltest)
            if pdf_path:
                document_builder.create_document(selected_questions, pdf_path, settings.ui.title_edit.text(),
                                                 icon=icon, font_size=settings.ui.fontsize_spinBox.value())
                webbrowser.open_new(pdf_path)

            if csv_path:
                with open(csv_path, 'w+', newline='', encoding='utf-8') as file:
                    # create the csv writer
                    writer = csv.writer(file)

                    for question in selected_questions:
                        # write a row to the csv file
                        writer.writerow([question.question.question.replace("\n", "\\n"),
                                         question.question.answer_text.replace("\n", "\\n")])
            if ppt_path:
                prs = pptx.Presentation(os.path.join(base_path, 'res/template.pptx'))
                questions_a = list(selected_questions)
                num_groups = settings.ui.spinBox_ppt_groups.value()

                def delete_slide(i):
                    rId = prs.slides._sldIdLst[i].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]

                delete_slide(1)
                delete_slide(0)

                slide_name = "Regelfrage_1"
                group_info = "Eine Gruppe - A"
                group_detailed_info = []

                if num_groups == 2:
                    questions_a, questions_b = generate_powerpoint_questionsets(questions_a)
                    slide_name = "Regelfrage_2"
                    group_info = "Zwei Gruppen - A und B"
                    group_detailed_info = [
                        "Gruppe A oben - blaue Schrift",
                        "Gruppe B unten - schwarze Schrift"
                    ]

                slide_layout = prs.slide_layouts.get_by_name(slide_name)

                slide = prs.slides.add_slide(prs.slide_layouts.get_by_name('Titelfolie'))
                slide.shapes.title.text = settings.ui.title_edit.text()

                slide = prs.slides.add_slide(prs.slide_layouts.get_by_name('1_Titel und Inhalt'))
                slide.shapes.title.text = "Aufbau"

                introduction_lines = [
                    f"{len(selected_questions)} Fragen",
                    group_info,
                    group_detailed_info,
                    f"{settings.ui.spinBox_ppt_time.value()} Sekunden Zeit pro Frage",
                    "Spielfortsetzung, persönliche Strafe, ggf. Ort",
                    "Einzelarbeit!"
                ]

                for shape in list(slide.shapes)[1:]:
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame
                    text_frame.clear()

                    p = text_frame.add_paragraph()
                    for line in introduction_lines:
                        if type(line) is list:
                            for sub_line in line:
                                p.level = 1
                                p.text = sub_line
                                p = text_frame.add_paragraph()
                        else:
                            p.level = 0
                            p.text = line
                            p = text_frame.add_paragraph()

                for i in range(len(questions_a)):
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.title.text = f"Frage {i + 1}"
                    list(slide.placeholders)[1].text = questions_a[i].question.question
                    if num_groups == 2:
                        list(slide.placeholders)[2].text = questions_b[i].question.question
                    slide._element

                prs.save(ppt_path)

                with open(str(ppt_path) + ".txt", 'w+') as file:
                    file.writelines("Lösungen Gruppe A\n")
                    file.writelines(
                        [f"Frage {i + 1}: {question.question.answer_text}\n" for i, question in enumerate(questions_a)])

                    if num_groups == 2:
                        file.writelines(["\n", "\n", "Lösungen Gruppe B\n"])
                        file.writelines([f"Frage {i + 1}: {question.question.answer_text}\n" for i, question in
                                         enumerate(questions_b)])

            QApplication.restoreOverrideCursor()


class SelfTestDockWidget(QWidget, Ui_self_test_dockwidget):
    changed = Signal()
    mode = SelfTestMode(0)  # type: SelfTestMode
    timer_question = Signal(int)
    timer_answer = Signal(int)

    def __init__(self, main_window: MainWindow):
        super(SelfTestDockWidget, self).__init__(main_window)
        self.ui = Ui_self_test_dockwidget()
        self.ui.setupUi(self)

        self.ui.self_test_question_groups.clear()
        self.ui.self_test_question_groups.setSelectionMode(QListWidget.ExtendedSelection)

        self.ui.question_visibility_spinbox.valueChanged.connect(self.timer_question)
        self.ui.auto_evaluate_spinbox.valueChanged.connect(self.timer_answer)

        self._question_groups = db.get_all_question_groups()
        for question in self._question_groups:
            item = QListWidgetItem(f"{question.id:02d} - {question.name}")
            item.setCheckState(Qt.Unchecked)
            self.ui.self_test_question_groups.addItem(item)

        self.ui.self_test_question_groups.itemChanged.connect(self._checkbox_changed)

        for mode in SelfTestMode:
            self.ui.mode_comboBox.addItem(str(mode))

        self.ui.mode_comboBox.currentIndexChanged.connect(self._combobox_changed)

    def lock(self):
        self.ui.mode_comboBox.setDisabled(True)
        self.ui.self_test_question_groups.setDisabled(True)

    def unlock(self):
        self.ui.mode_comboBox.setDisabled(False)
        self.ui.self_test_question_groups.setDisabled(False)

    def get_question_groups(self):
        question_groups = []
        for i, group in enumerate(self._question_groups):
            if self.ui.self_test_question_groups.item(i).checkState() == Qt.Checked:
                question_groups += [group]
        return question_groups

    def _checkbox_changed(self, item: QListWidgetItem):
        signal_state = self.ui.self_test_question_groups.blockSignals(True)
        new_state = item.checkState()
        selected_items = self.ui.self_test_question_groups.selectedItems()
        if item in selected_items:
            # otherwise just change the clicked item
            for item in selected_items:
                item.setCheckState(new_state)
        self.changed.emit()
        self.ui.self_test_question_groups.blockSignals(signal_state)

    def _combobox_changed(self, value: int):
        self.mode = SelfTestMode(value)
        self.changed.emit()

    def reset(self):
        signal_state = self.ui.self_test_question_groups.blockSignals(True)
        for index in range(self.ui.self_test_question_groups.count()):
            self.ui.self_test_question_groups.item(index).setCheckState(Qt.Unchecked)
        self.changed.emit()
        self.ui.self_test_question_groups.blockSignals(signal_state)
        self.ui.auto_evaluate_spinbox.setValue(0)
        self.ui.question_visibility_spinbox.setValue(0)
        self.ui.mode_comboBox.setCurrentIndex(0)
